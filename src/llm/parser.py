"""Extract text content from various file formats for LLM consumption."""

from __future__ import annotations

import logging
from pathlib import Path

logger = logging.getLogger(__name__)

# Maximum characters to extract from a single file (safety limit)
MAX_CHARS_PER_FILE = 100_000


def extract_text(file_path: str | Path) -> str:
    """Extract text from a file based on its extension.

    Supports: PDF, DOCX, PPTX, HTML, and plain text files.
    Returns empty string on failure.
    """
    path = Path(file_path)
    if not path.exists():
        return ""

    suffix = path.suffix.lower()

    try:
        if suffix == ".pdf":
            return _extract_pdf(path)
        elif suffix == ".docx":
            return _extract_docx(path)
        elif suffix in (".pptx", ".ppt"):
            return _extract_pptx(path)
        elif suffix in (".html", ".htm"):
            return _extract_html(path)
        elif suffix in _PLAIN_TEXT_EXTENSIONS:
            return _extract_plain(path)
        else:
            # Best effort: try as plain text
            return _extract_plain(path)
    except Exception as exc:
        logger.warning("Failed to extract text from %s: %s", path.name, exc)
        return ""


# Extensions we treat as plain text
_PLAIN_TEXT_EXTENSIONS = {
    ".txt", ".md", ".markdown", ".rst", ".csv", ".tsv", ".json",
    ".py", ".java", ".c", ".cpp", ".h", ".hpp", ".js", ".ts",
    ".rb", ".go", ".rs", ".sql", ".sh", ".bash", ".zsh",
    ".yaml", ".yml", ".toml", ".ini", ".cfg", ".conf",
    ".tex", ".bib", ".r", ".R", ".m", ".scala", ".kt",
    ".xml", ".css", ".scss", ".less", ".log",
}


def _extract_pdf(path: Path) -> str:
    """Extract text from PDF using PyMuPDF."""
    import fitz  # PyMuPDF

    text_parts: list[str] = []
    total = 0

    with fitz.open(str(path)) as doc:
        for page in doc:
            page_text = page.get_text()
            text_parts.append(page_text)
            total += len(page_text)
            if total > MAX_CHARS_PER_FILE:
                break

    return "\n".join(text_parts)[:MAX_CHARS_PER_FILE]


def _extract_docx(path: Path) -> str:
    """Extract text from DOCX."""
    from docx import Document

    doc = Document(str(path))
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
    return "\n".join(paragraphs)[:MAX_CHARS_PER_FILE]


def _extract_pptx(path: Path) -> str:
    """Extract text from PPTX (slide by slide)."""
    from pptx import Presentation

    prs = Presentation(str(path))
    text_parts: list[str] = []

    for i, slide in enumerate(prs.slides, 1):
        slide_texts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        slide_texts.append(text)
        if slide_texts:
            text_parts.append(f"[Slide {i}]")
            text_parts.extend(slide_texts)

    return "\n".join(text_parts)[:MAX_CHARS_PER_FILE]


def _extract_html(path: Path) -> str:
    """Extract text from HTML."""
    from bs4 import BeautifulSoup

    raw = path.read_text(encoding="utf-8", errors="replace")
    soup = BeautifulSoup(raw, "html.parser")

    # Remove script and style elements
    for tag in soup(["script", "style"]):
        tag.decompose()

    return soup.get_text(separator="\n", strip=True)[:MAX_CHARS_PER_FILE]


def _extract_plain(path: Path) -> str:
    """Read a plain text file."""
    try:
        return path.read_text(encoding="utf-8", errors="replace")[:MAX_CHARS_PER_FILE]
    except Exception:
        return ""


def chunk_text(text: str, max_tokens: int = 4000, overlap: int = 200) -> list[str]:
    """Split text into overlapping chunks by approximate token count.

    Uses a rough 1 token ≈ 4 chars heuristic for fast chunking.
    """
    if not text:
        return []

    char_limit = max_tokens * 4
    overlap_chars = overlap * 4

    chunks: list[str] = []
    start = 0
    while start < len(text):
        end = start + char_limit
        chunk = text[start:end]
        if chunk.strip():
            chunks.append(chunk.strip())
        start = end - overlap_chars

    return chunks


def extract_and_chunk(file_path: str | Path, max_tokens: int = 4000) -> list[str]:
    """Extract text from a file and split into chunks."""
    text = extract_text(file_path)
    if not text.strip():
        return []
    return chunk_text(text, max_tokens=max_tokens)
